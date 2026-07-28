"""Build the final-presentation .pptx from the guide content.

Everything here is generated: the two schematics on slides 2 and 4 are drawn as
native PowerPoint shapes (so they stay editable and count as original figures),
and the data charts come from results/figures/slides/*.png, which are re-plotted
from the committed run summaries at slide geometry.

The speaker script for each slide is written into the slide's notes, so the deck
and the script travel together.

Run:  .venv/bin/python scripts/build_deck.py
Out:  results/presentation/final-presentation.pptx
"""
from __future__ import annotations

import json
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
FIGS = os.path.join(ROOT, "results", "figures", "slides")
OUT = os.path.join(ROOT, "results", "presentation")

# one accent + two neutrals, per the rubric's "consistent color theme"
ACCENT = RGBColor(0x0F, 0x76, 0x6E)      # teal -- our contribution
INK = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x5A, 0x66, 0x72)
LIGHT = RGBColor(0xE8, 0xEE, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0x9A, 0x33, 0x24)         # same red as B0 in every chart

W, H = 13.333, 7.5


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, *, size, bold=False, color=INK, space_before=0, first=False,
         level=0, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(0)
    p.level = level          # python-pptx indents by outline level, not by inches
    f = p.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = "Calibri"
    return p


SLIDE_NO = 0


def chrome(slide, title, cite):
    """Title, accent rule, citation footer, slide number -- identical on every slide."""
    tf = textbox(slide, 0.55, 0.30, 12.2, 0.75)
    para(tf, title, size=25, bold=True, color=INK, first=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.02),
                                 Inches(2.1), Inches(0.055))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    if cite:
        tf = textbox(slide, 0.55, 6.92, 11.3, 0.4)
        para(tf, cite, size=10.5, color=MUTED, first=True)
    tf = textbox(slide, 12.35, 6.92, 0.6, 0.4)
    para(tf, str(SLIDE_NO), size=10.5, color=MUTED, first=True, align=PP_ALIGN.RIGHT)


def bullets(slide, blocks, x, y, w, *, main=17, sub=13.5, gap=9, h=5.4):
    """blocks = [(main_text, [sub, sub, ...]), ...]

    `h` is explicit because a text block sitting above a chart must be bounded --
    the default 5.4 in box would otherwise reach down over the figure.
    """
    tf = textbox(slide, x, y, w, h)
    first = True
    for head, subs in blocks:
        para(tf, head, size=main, bold=True, color=ACCENT,
             space_before=0 if first else gap, first=first)
        first = False
        for s in subs:
            para(tf, "–  " + s, size=sub, color=INK, level=1, space_before=2)


def note(slide, script):
    slide.notes_slide.notes_text_frame.text = script


MANIFEST = []                            # slide index -> figure file, for build_guide.py


def picture(slide, name, x, y, w):
    pic = slide.shapes.add_picture(os.path.join(FIGS, name), Inches(x), Inches(y),
                                   width=Inches(w))
    MANIFEST.append({"figure": name, "x": x, "y": y,
                     "w": round(w, 2), "h": round(pic.height / 914400.0, 2)})


def blank(prs):
    global SLIDE_NO
    SLIDE_NO += 1
    return prs.slides.add_slide(prs.slide_layouts[6])


# --------------------------------------------------------------- schematics
def desk_graphic(slide, x, y, w, h):
    """The exam-desk picture for slide 2, drawn natively so it stays editable."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = MUTED
    box.line.width = Pt(1.25)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    para(tf, "GPU memory  =  the desk", size=13, bold=True, color=MUTED, first=True,
         align=PP_ALIGN.CENTER)

    # eight "cheat sheets" -- the last one is the victim that gets preempted
    cols, cw, ch = 4, 0.82, 0.62
    gx, gy = x + 0.42, y + 0.62
    for i in range(8):
        r, c = divmod(i, cols)
        victim = i == 7
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(gx + c * (cw + 0.22)), Inches(gy + r * (ch + 0.26)),
                                   Inches(cw), Inches(ch))
        s.fill.solid()
        s.fill.fore_color.rgb = RED if victim else WHITE
        s.line.color.rgb = RED if victim else MUTED
        s.line.width = Pt(1.5 if victim else 1.0)
        s.shadow.inherit = False
        t = s.text_frame
        para(t, "cheat\nsheet", size=9, bold=victim,
             color=WHITE if victim else MUTED, first=True, align=PP_ALIGN.CENTER)
        t.vertical_anchor = MSO_ANCHOR.MIDDLE

    tf = textbox(slide, x, y + h + 0.04, w, 0.62)
    para(tf, "desk full  →  a student leaves  =  PREEMPTION", size=12.5,
         bold=True, color=RED, first=True, align=PP_ALIGN.CENTER)
    para(tf, "swap = kept in a locker  ·  recompute = thrown away",
         size=11, color=MUTED, align=PP_ALIGN.CENTER)


def arch_graphic(slide, x, y, w):
    """The layer diagram for slide 4: scheduler over our cache layer over the pool."""
    rows = [
        ("Requests", WHITE, MUTED, 0.52),
        ("LTR scheduler   (decides who runs next)", WHITE, MUTED, 0.62),
        ("KV control layer\nC1: 4-bit TurboQuant  ·  C2: sparsity τ = 6", LIGHT, ACCENT, 0.86),
        ("Paged KV cache pool on the GPU", WHITE, MUTED, 0.62),
    ]
    cy = y
    for i, (label, fill, line, hh) in enumerate(rows):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(cy),
                                   Inches(w), Inches(hh))
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        s.line.color.rgb = line
        s.line.width = Pt(2.25 if i == 2 else 1.0)
        s.shadow.inherit = False
        tf = s.text_frame
        tf.word_wrap = True
        para(tf, label, size=12.5 if i != 2 else 13, bold=(i == 2),
             color=ACCENT if i == 2 else INK, first=True, align=PP_ALIGN.CENTER)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if i < len(rows) - 1:
            a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                       Inches(x + w / 2 - 0.09), Inches(cy + hh + 0.03),
                                       Inches(0.18), Inches(0.24))
            a.fill.solid()
            a.fill.fore_color.rgb = MUTED
            a.line.fill.background()
            a.shadow.inherit = False
        cy += hh + 0.30

    tf = textbox(slide, x, cy - 0.06, w, 0.52)
    para(tf, "installed by a runtime hook — no engine patching", size=11.5,
         bold=True, color=ACCENT, first=True, align=PP_ALIGN.CENTER)



# ---------------------------------------------------------- extra primitives
def bullets2(slide, blocks, x, y, w, *, main=15.5, sub=12, gap=8, h=2.1, split=None):
    """Two-column bullet band, used on slides whose figure spans the full width."""
    n = split if split is not None else (len(blocks) + 1) // 2
    cw = (w - 0.45) / 2
    bullets(slide, blocks[:n], x, y, cw, main=main, sub=sub, gap=gap, h=h)
    bullets(slide, blocks[n:], x + cw + 0.45, y, cw, main=main, sub=sub, gap=gap, h=h)


def strip(slide, text, x, y, w, *, h=0.62, size=13.5):
    """A tinted callout bar -- the one sentence we want the room to remember."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = LIGHT
    s.line.color.rgb = ACCENT
    s.line.width = Pt(1.5)
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = True
    para(tf, text, size=size, bold=True, color=ACCENT, first=True, align=PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def grid(slide, rows, x, y, w, hh, *, size=11, head=11):
    """A small table. python-pptx tables carry a heavy default style, so every
    cell is restyled: white body, tinted header, no banding."""
    shape = slide.shapes.add_table(len(rows), len(rows[0]),
                                   Inches(x), Inches(y), Inches(w), Inches(hh))
    tbl = shape.table
    tbl.first_row = True
    tbl.horz_banding = False
    for r, row in enumerate(rows):
        tbl.rows[r].height = Inches(hh / len(rows))
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if r == 0 else WHITE
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            para(tf, val, size=head if r == 0 else size, bold=(r == 0 or c == 0),
                 color=ACCENT if r == 0 else INK, first=True)
    return tbl


# --------------------------------------------------------------------- deck
def build() -> str:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)

    # ---- 1  title -------------------------------------------------------
    s = blank(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(0.28))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

    tf = textbox(s, 0.62, 1.50, 12.1, 2.60)
    para(tf, "KVCache-Coordinated Latency Optimization", size=34, bold=True,
         color=INK, first=True)
    para(tf, "for LLM Inference Serving", size=34, bold=True, color=INK)
    para(tf, "Building on LTR scheduling with 4-bit KV compression and decode sparsity",
         size=18, color=MUTED, space_before=14)

    strip(s, "We evaluate KV-cache quantization and decode sparsity as a coordinated "
             "serving layer beneath an LTR scheduler, and separate memory-capacity "
             "gains from decode-kernel costs.",
          0.62, 4.32, 12.1, h=0.82, size=14)

    tf = textbox(s, 0.62, 5.42, 12.1, 1.5)
    para(tf, "Guoliang Liu   ·   Wenhui Kang   ·   Junpeng Huang", size=19,
         bold=True, color=ACCENT, first=True)
    para(tf, "Fairleigh Dickinson University — Vancouver Campus", size=15,
         color=INK, space_before=8)
    para(tf, "CSCI 6806 Graduate Capstone  ·  Group 3  ·  July 29, 2026", size=13,
         color=MUTED, space_before=4)
    note(s, """[0:15 — Guoliang]

Good afternoon. Our project is KVCache-Coordinated Latency Optimization for LLM
Inference Serving. At the midterm we proposed a KV-cache layer beneath a
learning-to-rank scheduler. Today we present the implementation, the controlled
experiments, and what we measured.""")

    # ---- 2  problem -----------------------------------------------------
    s = blank(prs)
    chrome(s, "The bottleneck: scheduling decides order, not memory cost",
           "[1] Kwon et al., PagedAttention, SOSP 2023  ·  [2] Fu et al., "
           "Learning-to-Rank scheduling, NeurIPS 2024")
    bullets(s, [
        ("Every decode step re-reads the whole KV cache", [
            "The cache grows with the prompt and with each token generated",
            "All concurrent requests share one paged pool on the GPU",
            "Memory, not arithmetic, becomes the admission bottleneck"]),
        ("When the pool fills, vLLM must evict a request", [
            "Swap: copy its blocks to host RAM and fetch them back later",
            "Recompute: discard its blocks and re-run prefill from scratch",
            "Either route adds a large, uneven penalty to that request"]),
        ("LTR fixes the ordering — it cannot shrink a request", [
            "An OPT-125M ranker predicts which outputs will be short",
            "Priority scheduling then reduces head-of-line blocking",
            "But a long request still occupies exactly as many blocks"]),
    ], 0.55, 1.40, 7.1, main=17, sub=13.5, gap=11, h=5.2)
    desk_graphic(s, 8.15, 1.75, 4.65, 2.55)
    tf = textbox(s, 8.15, 5.05, 4.65, 1.7)
    para(tf, "TTFT  time to the first token — dominated by queueing",
         size=12, color=INK, first=True)
    para(tf, "TPOT  average time per later token — dominated by the decode kernel",
         size=12, color=INK, space_before=6)
    para(tf, "Preemption  a running request evicted because the pool is full",
         size=12, color=RED, space_before=6)
    note(s, """[0:55 — Guoliang]

Serving a large model has two bottlenecks, not one. Every decode step re-reads
the KV cache of everything generated so far. That cache grows with the prompt and
with every token, and all concurrent requests share one pool of GPU memory.

When the pool fills, vLLM evicts a running request — either swapping its blocks to
host memory, or discarding them and re-running prefill later. Both are expensive,
and both land on one unlucky user.

Think of GPU memory as an exam desk. Each student is a request; their cheat sheet
is the KV cache. Desk full, somebody leaves — that is a preemption.

Learning-to-rank decides who sits down first, which reduces head-of-line
blocking. What it cannot do is make any cheat sheet smaller. That gap is what we
went after.""")

    # ---- 3  where the cost lives (report Fig. 1) ------------------------
    s = blank(prs)
    chrome(s, "Where the cost lives: the cache and its escape routes",
           "Fig. 1 of our report  ·  [1] Kwon et al., 2023")
    bullets(s, [
        ("The serving path we instrument end to end", [
            "Arrivals → LTR ranker → priority queue → paged block allocator",
            "Each decode step reads K and V blocks for every running request",
            "We log per-request TTFT, TPOT and E2E, plus per-run preemptions"]),
        ("Preemption is the measurable symptom of memory pressure", [
            "Engine A (vLLM 0.4.1 fork) recovers by swapping blocks to host RAM",
            "Engines B and C (0.8.5, 0.25.0) recover by recomputing the prefill",
            "We measured both, so no finding rests on a single recovery policy"]),
        ("Our layer changes the byte cost, not the allocator", [
            "C1 changes how many bytes one cache token occupies",
            "The scheduler and the block allocator are left untouched",
            "So a gain is attributable to capacity, not to a new policy"]),
    ], 0.55, 1.40, 7.5, main=16.5, sub=13, gap=11, h=5.2)
    picture(s, "slide_arch_preempt.png", 8.35, 1.62, 4.45)
    note(s, """[0:42 — Guoliang]

This is the path a request actually takes, and the path we instrumented.
Arrivals, the ranker, the priority queue, the paged allocator. Every decode step
reads keys and values for every running request.

The part that matters is on the right: what happens when the allocator runs out.
Older vLLM swaps to host memory; newer vLLM deletes the blocks and re-runs
prefill. We measured both, so nothing here depends on one recovery policy.

And note where our layer sits. We do not touch the scheduler or the allocator. We
change one thing — how many bytes a cache token costs.

SHORT PATH (0:12) — say only: "This is the serving path we instrumented. The part that matters is what happens
when the allocator runs out: older vLLM swaps to host memory, newer vLLM re-runs
prefill. We measured both."
""")

    # ---- 4  where sparsity attaches (report Fig. 2) ---------------------
    s = blank(prs)
    chrome(s, "Where sparsity attaches: the decode attention kernel",
           "Fig. 2 of our report  ·  [6] Yuan et al., BLASST, 2025")
    bullets(s, [
        ("Decode attention is a streaming softmax over cache blocks", [
            "The kernel walks the KV blocks keeping a running max and sum",
            "Each block costs a K read, a V read, and an accumulate",
            "Most of the time goes into moving V, not into computing scores"]),
        ("BLASST gates on the running max, before the V read", [
            "A block whose best score falls below m − τ contributes ≈ 0 weight",
            "The kernel then skips that block's V load and its accumulation",
            "The K read still happens, so the gate itself is not free"]),
        ("A decode-side change only", [
            "Prefill, cache layout and sampling are all unchanged",
            "τ = 6 was chosen from a perplexity sweep and then frozen",
            "Output quality is verified afterwards on GSM8K, not assumed"]),
    ], 0.55, 1.40, 7.5, main=16.5, sub=13, gap=11, h=5.2)
    picture(s, "slide_arch_attention.png", 8.35, 1.58, 4.45)
    note(s, """[0:50 — Guoliang]

The second lever lives inside the attention kernel. Decode attention is a
streaming softmax: the kernel walks the cache blocks keeping a running maximum
and sum, and for each block it reads keys, reads values, and accumulates. Most of
the time goes into moving values, not computing scores.

BLASST adds a gate after the scores and before the values are read. If a block's
best score is more than tau below the running maximum — we use six — its weight
after the softmax is negligible, so the kernel skips loading its values.

Two caveats: the keys are still read, so the gate is not free, and tau came from a
perplexity sweep we froze before these runs.

SHORT PATH (0:15) — say only: "Decode attention streams over the cache blocks. BLASST skips a block's value load
when its score is far below the running maximum. Our threshold is six, frozen
from a perplexity sweep."
""")

    # ---- 5  what we built ----------------------------------------------
    s = blank(prs)
    chrome(s, "What we built: a coordinated KV layer under LTR",
           "[2] Fu et al., 2024  ·  [4] TurboQuant, 2025  ·  [6] BLASST, 2025")
    bullets(s, [
        ("We keep LTR and add two levers underneath it", [
            "LTR still answers: which request runs next?",
            "C1 answers: how many bytes does one cache token cost?",
            "C2 answers: which blocks can the decode kernel skip?"]),
        ("C1 — TurboQuant 4-bit KV cache, native in vLLM 0.25.0", [
            "Cache entries stored at 4 bits instead of bf16",
            "2.83× more cache tokens fit in the same GPU bytes",
            "Cost: every decode step has to unpack the low-bit values"]),
        ("C2 — BLASST sparsity, installed by a runtime hook", [
            "A sitecustomize.py on PYTHONPATH rewrites the compiled decode kernel",
            "The hook reaches vLLM's spawned EngineCore workers, where it runs",
            "The installed vLLM package is never edited — reversible and auditable"]),
    ], 0.55, 1.40, 7.5, main=16.5, sub=13, gap=11, h=4.5)
    arch_graphic(s, 8.45, 1.55, 4.35)
    strip(s, "Contribution: not “latency improved”, but a coordinated layer in which "
             "memory-capacity gains and decode-kernel costs are measured separately.",
          0.55, 5.93, 7.5, h=0.84, size=13)
    note(s, """[1:09 — Guoliang]

Here is what we actually built. We kept the learning-to-rank scheduler and added
two levers underneath it. Three questions: LTR decides which request runs next,
C1 decides how many bytes a cache token costs, C2 decides which blocks the kernel
can skip.

C1 is TurboQuant four-bit KV, using vLLM 0.25's native support. At equal GPU
memory it fits 2.83 times more cache tokens — the same cheat sheet in smaller
handwriting. But smaller handwriting is slower to read, and that shows up as
dequantization cost on every decode step.

C2 is where most of our engineering went. BLASST ships in no released vLLM, and
we did not want to fork the engine. So we wrote a sitecustomize module on
PYTHONPATH that rewrites the compiled decode kernel in memory. The subtlety is
that vLLM spawns separate EngineCore worker processes, so the hook has to reach
into those workers. The installed package is never modified — the patch is
reversible, and a reviewer can diff it.""")

    # ---- 6  platform (report Table I) -----------------------------------
    s = blank(prs)
    chrome(s, "Platform, model and workload",
           "Table I of our report  ·  every row is held fixed across all runs")
    bullets(s, [
        ("One consumer GPU, three engine generations", [
            "RTX 3090, 24 GB, SM86 — not a datacenter card",
            "A: vLLM 0.4.1 fork (swap) · B: 0.8.5 V0 · C: 0.25.0 V1",
            "Engine C is the one with native TurboQuant KV support"]),
        ("A real chat trace, not synthetic prompts", [
            "LMSYS-Chat-1M for the main sweep, ShareGPT for the reproduction",
            "Poisson arrivals, 4–64 req/s, n = 200 per rate on Engine C",
            "The ranker is ours: OPT-125M, ListMLE, 23,800 samples"]),
        ("Everything needed to repeat the runs is pinned", [
            "Llama-3.1-8B-Instruct, bf16 (fp16 on Engine B)",
            "Python 3.10, torch 2.11, Triton 3.6, under WSL2",
            "GSM8K for quality, plus the perplexity sweep that fixed τ"]),
    ], 0.55, 1.40, 6.6, main=16, sub=12.5, gap=10, h=5.3)
    grid(s, [
        ["Component", "Specification"],
        ["GPU", "RTX 3090, 24 GB, SM86 (Ampere)"],
        ["Host", "WSL2, Linux 6.6.87, 22 GB RAM"],
        ["Model", "Llama-3.1-8B-Instruct, bf16 (fp16 on B)"],
        ["Engine A", "vLLM 0.4.1 fork, swap preemption"],
        ["Engine B", "vLLM 0.8.5.post1, V0 engine, recompute"],
        ["Engine C", "vLLM 0.25.0, V1 engine, TurboQuant KV"],
        ["Runtime", "Python 3.10, torch 2.11, Triton 3.6"],
        ["Ranker", "OPT-125M, ListMLE, 23,800 samples, 10 epochs"],
        ["Workload", "LMSYS-Chat-1M; ShareGPT (reproduction)"],
        ["Load", "Poisson, 4–64 req/s (2–60 reproduction)"],
        ["Requests", "n = 200/rate on C, 100 on A and B"],
        ["Accuracy", "GSM8K, n = 100, greedy; perplexity"],
    ], 7.45, 1.42, 5.35, 5.05, size=10.5, head=11)
    note(s, """[0:35 — Wenhui]

Before the results, the platform — because on a project like this the hardware is
part of the claim.

One RTX 3090. A consumer card, not a datacenter GPU, and that shapes how we had
to design the experiment. Llama-3.1-8B-Instruct, and three vLLM generations: the
0.4.1 fork the LTR paper used, 0.8.5, and 0.25 which is the one with native
TurboQuant.

The workload is a real chat trace, LMSYS-Chat-1M, replayed with Poisson arrivals
from 4 to 64 requests per second, 200 requests at each rate. The reproduction
uses ShareGPT, because that is what the original work used. The ranker is ours:
OPT-125M trained with ListMLE on 23,800 samples.

Every row in this table is held fixed across every run in this talk, so anything
you see change later comes from the configuration, not from the setup.

SHORT PATH (0:12) — say only: "One RTX 3090, Llama-3.1-8B, three vLLM
generations, and a real chat trace replayed from 4 to 64 requests per second.
Every row here is held fixed across all our runs."
""")

    # ---- 7  experimental design ----------------------------------------
    s = blank(prs)
    chrome(s, "Experimental design: what makes the comparison fair",
           "Table II and Section IV of our report  ·  LMSYS-Chat-1M workload")
    bullets(s, [
        ("Equal GPU bytes, not equal block counts", [
            "bf16:  512 blocks  =  8,192 cache tokens",
            "TQ4:  724 blocks  =  23,168 cache tokens",
            "The two pools occupy approximately the same device memory"]),
        ("A five-rung ladder, one change per rung", [
            "Each rung differs from the one below it by exactly one setting",
            "Same trace, same seed, same arrival process on every rung",
            "Poisson arrivals, 4–64 req/s, n = 200 requests per rate"]),
        ("Two controls separate capacity from codec", [
            "ctrl: bf16 given TQ4's token capacity → capacity effect alone",
            "c1te: TQ4 held to the baseline token budget → 4-bit cost alone",
            "Without these, “compression helps” and “more room helps” are confounded"]),
        ("The pool cap is a stated design decision", [
            "An uncapped RTX 3090 goes compute-limited before the cache fills",
            "We cap it on purpose to reach the memory-bound regime we study"]),
    ], 0.55, 1.40, 6.85, main=16, sub=12.5, gap=9, h=5.3)
    grid(s, [
        ["Config", "Scheduler", "KV dtype", "Attention", "Pool (blk/tok)"],
        ["B0", "FCFS", "bf16", "dense", "512 / 8,192"],
        ["B1", "priority", "bf16", "dense", "512 / 8,192"],
        ["B1+C2", "priority", "bf16", "τ = 6", "512 / 8,192"],
        ["B1+C1", "priority", "TQ4", "dense", "724 / 23,168"],
        ["B1+C1+C2", "priority", "TQ4", "τ = 6", "724 / 23,168"],
        ["ctrl", "priority", "bf16", "dense", "1,448 / 23,168"],
        ["c1te", "priority", "TQ4", "dense", "256 / 8,192"],
    ], 7.75, 1.62, 5.05, 3.5, size=10.5)
    strip(s, "ctrl costs twice the memory to reach TQ4's capacity; c1te pays the 4-bit "
             "decode cost with none of the capacity. Together they split the two effects.",
          7.75, 5.42, 5.05, h=1.05, size=12)
    note(s, """[1:10 — Wenhui]

This is the slide I would most like you to hold us to, because it is where a
result like ours is usually unfair.

The obvious mistake is to give four-bit more cache blocks and then call it fast.
So we equalise bytes, not blocks: bf16 gets 512 blocks, which is 8,192 cache
tokens; four-bit gets 724 blocks, which is 23,168. Same GPU memory. The extra
tokens are exactly what compression buys.

Second, every rung of the ladder differs from the one below it by one setting, on
the same trace with the same seed. And the two controls at the bottom: "ctrl"
gives bf16 four-bit's capacity, isolating capacity alone; "c1te" holds four-bit
to the baseline budget, isolating the codec cost. Without them, "compression
helps" and "more room helps" are the same sentence.

One thing before you ask: our 3090 is small, and uncapped it goes compute-limited
before the cache fills. Capping the pool is how we reach the memory-bound regime
this work is about.""")

    # ---- 8  baseline reproduction (report Fig. 3) ----------------------
    s = blank(prs)
    chrome(s, "We earned the baseline first: reproducing the LTR result",
           "Fig. 3 of our report  ·  [2] Fu et al., 2024  ·  [3] Kumar et al., 2026")
    bullets(s, [
        ("We reproduced the published LTR result before adding anything", [
            "Engine A — the vLLM 0.4.1 fork with swap preemption",
            "ShareGPT trace, B0 first-come-first-served vs B1 LTR priority",
            "The original paper's metric: normalized latency"]),
        ("The benefit appears only once memory is under pressure", [
            "At low arrival rates B0 and B1 are indistinguishable",
            "The curves separate only after the block pool saturates",
            "The same regime dependence returns later for compression"]),
        ("Our measurements bracket the published number", [
            "Mean normalized latency: 1.66× better",
            "P99 normalized latency: 3.6× better",
            "The published 2.1× sits between our mean and our tail"]),
    ], 0.55, 1.40, 7.1, main=16.5, sub=13, gap=11, h=5.2)
    picture(s, "slide_repro.png", 7.95, 2.35, 4.9)
    note(s, """[0:42 — Wenhui]

Before adding anything of our own, we reproduced the baseline we build on — an
assumed baseline makes every improvement above it unverifiable.

Engine A, the 0.4.1 fork with swap preemption, ShareGPT trace,
first-come-first-served against LTR priority, the original paper's metric.

Two results. The shape: at low rates the two are indistinguishable, because the
pool still has free blocks; they separate only under memory pressure — and
compression behaves the same way later. The magnitude: 1.66 times on the mean,
3.6 times at the 99th percentile. The published 2.1 sits between them. So we
treat LTR as an earned baseline.""")

    # ---- 9  headline (report Fig. 5) -----------------------------------
    s = blank(prs)
    chrome(s, "Headline: each lever fixes a different bottleneck",
           "Fig. 5 of our report  ·  Engine C (vLLM 0.25.0), 64 req/s, n = 200 per rate")
    bullets2(s, [
        ("B0 → B1   scheduling drains the queue", [
            "Mean TTFT  12,275 → 6,488 ms",
            "Mean TPOT  80 → 68 ms/token",
            "Mean E2E  20,003 → 15,063 ms"]),
        ("B1 → B1+C1   compression removes the pressure", [
            "Mean TTFT  6,488 → 2,347 ms",
            "But mean TPOT  68 → 215 ms/token — the 4-bit decode tax",
            "E2E rises to 20,708 ms: capacity alone is not a win"]),
        ("+C2   sparsity refunds part of that tax", [
            "Mean TPOT  215 → 167 ms/token  (−22.1%)",
            "Mean TTFT  2,347 → 1,747 ms",
            "Mean E2E  20,708 → 17,308 ms — now below B1"]),
    ], 0.55, 1.32, 12.2, main=14.5, sub=11.5, gap=7, h=2.05, split=2)
    picture(s, "slide_headline.png", 1.57, 3.48, 10.2)
    note(s, """[0:58 — Wenhui]

The headline, at the hardest load we ran, 64 requests per second.

Left to right. Scheduling alone nearly halves mean time-to-first-token, 12.3
seconds to 6.5 — that is queueing delay draining. Four-bit compression cuts it
again, to 2.3 seconds. But look at the middle panel: time per output token goes
from 68 milliseconds to 215, and end-to-end latency goes back up, to 20.7
seconds.

That is the honest result. Compression is not free acceleration. It solves a
memory problem and charges you on every decode step, because low-bit values have
to be unpacked. Capacity by itself is not a win.

Now turn on sparsity. Per-token drops from 215 to 167, a 22 percent cut,
time-to-first-token to 1.7 seconds, and end-to-end lands at 17.3 — below
LTR-only. The two levers are complementary because they act on different
resources.""")

    # ---- 10  cross-engine (report Fig. 4) -------------------------------
    s = blank(prs)
    chrome(s, "The capacity gain is not an artefact of one engine",
           "Fig. 4 and Table II of our report  ·  three vLLM generations")
    bullets2(s, [
        ("Three engine generations, two recovery policies", [
            "A — vLLM 0.4.1 fork, swap: the LTR paper's own platform",
            "B — vLLM 0.8.5 V0, recompute",
            "C — vLLM 0.25.0 V1, recompute, native TurboQuant"]),
        ("The same qualitative ladder appears on all three", [
            "Under pressure, more cache capacity always lowers TTFT",
            "The absolute numbers differ; the ordering does not",
            "So the effect follows the memory regime, not the codebase"]),
        ("The residual differences are policy, not noise", [
            "Swap pays a host transfer; recompute pays a second prefill",
            "Occupancy 94.5% on the swap engine vs 100% on recompute",
            "Porting the workload across versions was a large part of the work"]),
    ], 0.55, 1.32, 12.2, main=14.5, sub=11.5, gap=7, h=2.05, split=2)
    picture(s, "slide_engines.png", 1.13, 3.48, 11.08)
    note(s, """[0:41 — Wenhui]

A fair question is whether this is an accident of one vLLM version, so we ran the
ladder on three.

Engine A is the fork the LTR paper itself used, which recovers by swapping. B and
C are much newer and recover by recomputing; C is the one with native TurboQuant.
Porting the same workload and the same metrics across three engine generations
was a substantial piece of work on its own.

The same qualitative ladder appears on all three. The absolute numbers differ;
the ordering does not. So the effect follows the memory regime, not the
codebase.

SHORT PATH (0:12) — say only: "We ran the same ladder on three vLLM generations with two different recovery
policies. The absolute numbers differ; the ordering is the same on all three."
""")

    # ---- 11  tail (report Fig. 6) --------------------------------------
    s = blank(prs)
    chrome(s, "The gain lands exactly where users feel it: the tail",
           "Fig. 6 of our report  ·  TTFT percentiles from the same runs")
    bullets(s, [
        ("We report the distribution, not only the mean", [
            "P50, P75, P90, P95 and P99, all from the same runs",
            "A mean hides what queueing does to the unlucky request",
            "Tail latency is what an interactive user actually notices"]),
        ("Compression flattens the tail hardest", [
            "The baseline curve climbs steeply past P75 — that climb is preemption",
            "The compressed curves stay flat much further out",
            "The gap between configurations widens as the percentile rises"]),
        ("The reason is mechanical, not statistical", [
            "A preempted request pays its queueing cost twice",
            "Preventing the eviction removes its whole contribution to the tail",
            "Which is why the tail improves more than the mean does"]),
    ], 0.55, 1.40, 7.1, main=16.5, sub=13, gap=11, h=5.2)
    picture(s, "slide_tail.png", 7.95, 2.35, 4.9)
    note(s, """[0:30 — Wenhui]

Means hide the thing users complain about, so we report the whole distribution.

Look at where the baseline curve turns upward — past the 75th percentile it climbs
steeply. That climb is preemption: a few requests paying a very large penalty,
which a mean quietly averages away. The compressed curves stay flat much further
out, because a preempted request pays its queueing cost twice, and preventing the
eviction removes that contribution entirely.

SHORT PATH (0:12) — say only: "The baseline's tail climbs steeply past P75, and that climb is preemption.
Compression flattens it, because a preempted request pays its queueing cost
twice."
""")

    # ---- 12  sign flip (report Fig. 7a) --------------------------------
    s = blank(prs)
    chrome(s, "Key finding: the same algorithm flips sign by kernel",
           "Fig. 7(a) of our report  ·  [6] BLASST, 2025  ·  identical τ = 6 in both arms")
    bullets(s, [
        ("Identical gate, identical threshold, opposite result", [
            "TQ4 per-head kernel:  mean TPOT −22.1%,  P99 TPOT −46.8%",
            "bf16 unified GQA kernel:  mean TPOT +4.9%,  P99 +15.0%",
            "Nothing changed except the kernel underneath"]),
        ("The cause is tile sharing in grouped-query attention", [
            "Four query heads share one KV tile in the unified kernel",
            "A tile can be skipped only if all four heads agree to skip it",
            "One dissenting head forces the whole tile to be processed"]),
        ("Identified sparsity is not realized sparsity", [
            "The gate identifies ≈ 42% of blocks as negligible",
            "The GQA kernel is only able to realize ≈ 18.5%",
            "Below break-even the gate costs more than the work it saves"]),
    ], 0.55, 1.40, 7.1, main=16.5, sub=13, gap=11, h=5.2)
    picture(s, "slide_signflip.png", 7.95, 1.62, 4.9)
    strip(s, "An algorithm can identify a block as skippable while the kernel remains "
             "physically unable to skip it.",
          7.95, 4.85, 4.9, h=0.85, size=12.5)
    note(s, """[1:03 — Junpeng]

This is the finding we think is worth more than any single latency number.

Same algorithm, same threshold, two places. On the four-bit per-head kernel it
cuts time per output token by 22 percent, and by 47 percent at the 99th
percentile. On the bf16 unified kernel the identical algorithm is 4.9 percent
worse on the mean and 15 percent worse at the tail. Nothing changed except the
kernel underneath.

The cause is grouped-query attention. Four query heads share one execution tile.
Head A may decide a block is irrelevant, but if head B still needs it, the tile
runs anyway — skipping requires unanimity. Our gate identifies about 42 percent of
blocks as negligible; the grouped kernel can skip only 18.5. Below break-even the
gate costs more than it saves.

The lesson: algorithmic sparsity is not automatically hardware speedup — the tile
structure of the kernel decides.""")

    # ---- 13  compute density (report Fig. 7b) --------------------------
    s = blank(prs)
    chrome(s, "Second condition: decode has to be compute-dense",
           "Fig. 7(b) of our report  ·  context sweep at a fixed 6 req/s")
    bullets(s, [
        ("We swept context length at a fixed arrival rate", [
            "Longer contexts mean fewer requests fit concurrently",
            "So context length and batch size move in opposite directions",
            "That lets us separate “long context” from “dense batch”"]),
        ("The benefit tracks batch size, not context length", [
            "512-token context, batch 40:  TPOT −20.4%",
            "2,048-token context, batch 21:  −10.3%",
            "4,096 (batch 10) and 7,168 (batch 6):  +0.6% and +2.8%"]),
        ("Because the gate is a fixed per-block overhead", [
            "Its cost is amortised across the requests sharing a batch",
            "A thin batch cannot amortise it, so the sign flips again",
            "Sparsity is a compute-side tool; it needs compute to be the bottleneck"]),
    ], 0.55, 1.40, 7.1, main=16.5, sub=13, gap=11, h=5.2)
    picture(s, "slide_density.png", 7.95, 2.35, 4.9)
    note(s, """[0:41 — Junpeng]

The kernel is one condition. The second is the workload.

We swept context length at a fixed arrival rate. The thing to notice is that
context length and batch size move in opposite directions — longer contexts mean
fewer requests fit concurrently — so this sweep separates "long context" from
"dense batch".

At 512 tokens, with a batch of 40, sparsity cuts per-token time by 20 percent. By
4,096 tokens the batch is 10 and the benefit is gone. The gate is a fixed
per-block cost, amortised across the batch, and a thin batch cannot amortise
it.

SHORT PATH (0:14) — say only: "The benefit tracks batch size, not context length. At batch 40 sparsity saves 20
percent; by batch 10 it is gone. The gate is a fixed per-block cost, and a thin
batch cannot amortise it."
""")

    # ---- 14  mechanism + accuracy (report Fig. 8) ----------------------
    s = blank(prs)
    chrome(s, "The mechanism is measured, not inferred",
           "Fig. 8 of our report  ·  preemptions and peak cache occupancy per run")
    bullets2(s, [
        ("Preemption counts follow the TTFT story exactly", [
            "B0 162  →  B1 147  →  B1+C1 42  →  B1+C1+C2 40 per run",
            "Scheduling barely moves it; compression collapses it",
            "So the TTFT gain has a named cause, not a correlation"]),
        ("Occupancy confirms the memory explanation", [
            "Uncompressed pools sit at 94.5% (swap) to 100% (recompute)",
            "Compressed pools settle at roughly 75.0–77.6%",
            "Swap reacts after overflow; compression prevents it"]),
        ("Swap and compression fight the same pressure, at different times", [
            "Swap reacts only once the pool has already overflowed",
            "Compression stops the pool reaching the overflow point",
            "That is the whole difference between 94.5% and 77.6% occupancy"]),
    ], 0.55, 1.32, 12.2, main=14.5, sub=11.5, gap=7, h=2.05, split=2)
    picture(s, "slide_preempt.png", 1.57, 3.48, 10.2)
    note(s, """[0:54 — Junpeng]

We did not want to stop at "it got faster", so we measured the mechanism.

The left panel is preemptions per run. First-come-first-served, 162. LTR, 147 —
scheduling barely moves it, because reordering does not create memory. Add
compression, 42. Add sparsity, 40. The collapse happens exactly where compression
switches on, so the time-to-first-token gain has a named cause.

Occupancy confirms it: uncompressed pools run 94.5 to 100 percent full,
compressed settle around 75 to 78. That is the cleanest way to say it — swapping
reacts after overflow, compression prevents the overflow.

And this is the cleanest statement of what compression does that swapping does
not: swapping reacts after the overflow has already happened, compression stops
the pool from getting there.""")

    # ---- 15  accuracy (report Table III) --------------------------------
    s = blank(prs)
    chrome(s, "Quality check: both levers are lossy — did it cost accuracy?",
           "Table III of our report  ·  GSM8K, n = 100, greedy decoding, identical prompts")
    bullets(s, [
        ("Both levers throw information away, so this had to be measured", [
            "C1 stores keys and values at 4 bits instead of 16",
            "C2 never loads the values of blocks the gate rules out",
            "Neither is lossless, so a latency win could hide a quality loss"]),
        ("No configuration falls below the bfloat16 baseline", [
            "The whole spread is +0.03 to +0.08, and every sign is positive",
            "τ = 6 was fixed by a perplexity sweep before any of these runs",
            "Same 100 problems, same greedy decoding, in every row"]),
        ("We read this as no measured degradation, not as an improvement", [
            "At n = 100 a swing of ±0.05 sits inside sampling noise",
            "C2 alone scores highest at 0.88 — noise explains it, so we do not claim it",
            "A real accuracy claim needs thousands of problems; we list that as a limit"]),
    ], 0.55, 1.40, 7.1, main=16, sub=12.5, gap=11, h=5.2)
    grid(s, [
        ["Configuration", "GSM8K", "Δ vs. bf16"],
        ["bfloat16, dense (baseline)", "0.80", "—"],
        ["bfloat16 + BLASST (C2)", "0.88", "+0.08"],
        ["TurboQuant-4bit (C1)", "0.83", "+0.03"],
        ["TQ4 + BLASST (C1+C2)", "0.85", "+0.05"],
    ], 7.95, 2.35, 4.9, 2.0, size=11.5)
    strip(s, "Lossy on paper, not lossy in practice — at the sample size we can afford "
             "to run.", 7.95, 4.75, 4.9, h=0.85, size=12.5)
    note(s, """[0:32 — Junpeng]

One more thing we owed you. Both of our levers are lossy: four-bit quantization
throws away precision, and sparsity skips blocks entirely. So the obvious worry
is that we bought latency with answer quality.

We measured it. GSM8K, a hundred problems, greedy decoding, the same prompts in
every row. Baseline bfloat16 scores 0.80. Four-bit scores 0.83. The full system
scores 0.85. Nothing falls below the baseline.

We want to be careful about how we read this. At a hundred problems, a swing of
five points is inside sampling noise, and sparsity alone actually scores highest
at 0.88 — which we do not believe means sparsity improves reasoning. The
defensible claim is that we found no measured degradation, not that our system is
more accurate. Getting a tighter number needs thousands of problems, and we list
that as a limitation.

SHORT PATH (0:12) — say only: "Both levers are lossy, so we checked. GSM8K at
a hundred problems: 0.80 baseline, 0.83 four-bit, 0.85 full system. Nothing below
baseline — but at that sample size the honest claim is no measured degradation,
not higher accuracy."
""")

    # ---- 16  conclusion -------------------------------------------------
    s = blank(prs)
    chrome(s, "Conclusion: coordinate capacity with decode efficiency",
           "Full results, limitations and reproduction scripts in our final report")
    bullets2(s, [
        ("What we built and verified", [
            "LTR scheduling + 4-bit TurboQuant KV + BLASST decode sparsity as one layer",
            "A runtime kernel patch that never edits the installed engine",
            "Reproduced baseline, equal-byte controls, three engines, latency and accuracy"]),
        ("What the measurements say", [
            "2.83× more cache tokens at equal device memory",
            "TTFT 12,275 → 1,747 ms and preemptions 162 → 40 at 64 req/s",
            "Compression pays when the server is memory-bound",
            "Sparsity pays only when the kernel can realize what the gate identifies"]),
        ("Limits and next steps", [
            "One RTX 3090, with the block pool capped on purpose",
            "A single fixed threshold; one production kernel patched",
            "Next: a per-head GQA kernel, an adaptive τ, a datacenter-GPU replication"]),
    ], 0.55, 1.38, 12.2, main=15.5, sub=12.5, gap=9, h=3.5, split=2)
    strip(s, "Scheduling controls order  ·  compression controls footprint  ·  "
             "sparsity controls decode cost — and each one only pays in its own regime.",
          0.55, 5.35, 12.2, h=0.78, size=15)
    tf = textbox(s, 0.55, 6.28, 12.2, 0.5)
    para(tf, "Thank you — questions welcome.", size=15, bold=True, color=INK,
         first=True, align=PP_ALIGN.CENTER)
    note(s, """[1:03 — Junpeng]

To close.

We built a coordinated KV layer beneath a learning-to-rank scheduler: four-bit
TurboQuant for capacity, BLASST sparsity for decode cost, installed by a runtime
patch that never edits the engine. And we tested it the way we would want to be
tested: reproduce the baseline, equalise memory, run controls, check accuracy as
well as latency.

Our conclusion is conditional. Compression pays when the server is memory-bound.
Sparsity pays only when the kernel can realize what the gate identifies, and the
batch is dense enough to amortise it.

The limits are real — one 3090, a capped pool, one fixed threshold, one kernel —
and the next steps follow from our own findings: a per-head grouped-query kernel,
an adaptive threshold, a datacenter-GPU replication.

One sentence to take away: scheduling controls order, compression controls
footprint, sparsity controls decode cost, and each one only pays in its own
regime. Thank you.""")

    os.makedirs(OUT, exist_ok=True)
    path = os.path.join(OUT, "final-presentation.pptx")
    prs.save(path)
    with open(os.path.join(OUT, "deck-figures.json"), "w", encoding="utf-8") as fh:
        json.dump(MANIFEST, fh, indent=2)
    return path


if __name__ == "__main__":
    p = build()
    print("wrote", p)
