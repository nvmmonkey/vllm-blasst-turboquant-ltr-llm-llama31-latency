"""Generate docs/final-presentation-guide.md from the built .pptx.

The per-slide sections -- titles, bullets, citation lines, figure assignments,
click order and speaker script -- are read back out of the deck, so the guide
cannot drift away from the file you actually present. The narrative sections
(work split, analogy, rubric, Q&A) live here as text.

Run:  .venv/bin/python latex_source/scripts/build_guide.py
"""
from __future__ import annotations

import json
import os
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
DECK = os.path.join(ROOT, "results", "presentation", "final-presentation.pptx")
MANI = os.path.join(ROOT, "results", "presentation", "deck-figures.json")
OUT = os.path.join(ROOT, "docs", "final-presentation-guide.md")
EMU = 914400.0
ACCENT = "0F766E"          # RGBColor is not an int subclass; compare as hex text

# figure on the slide -> the figure it corresponds to in the report
REPORT_FIG = {
    "slide_arch_preempt.png": "Fig. 1 (arch_preempt)",
    "slide_arch_attention.png": "Fig. 2 (arch_attention)",
    "slide_repro.png": "Fig. 3 (fig7_nlatency_ladder)",
    "slide_engines.png": "Fig. 4 (fig1_cross_engine_ladder)",
    "slide_headline.png": "Fig. 5 (fig2_five_config_r64)",
    "slide_tail.png": "Fig. 6 (fig3_ttft_percentiles)",
    "slide_signflip.png": "Fig. 7a (fig4_c2_regimes)",
    "slide_density.png": "Fig. 7b (fig4_c2_regimes)",
    "slide_preempt.png": "Fig. 8 (fig6_preemptions)",
}
NATIVE = {
    2: "**Native PowerPoint shapes** — the exam-desk graphic, drawn by "
       "`desk_graphic()`. Editable, and original rather than a screenshot.",
    5: "**Native PowerPoint shapes** — the four-layer stack, drawn by "
       "`arch_graphic()`. Editable, and original rather than a screenshot.",
    6: "**Native table — Table I of the report** (experimental platform), "
       "drawn by `grid()`. Rows match the report cell for cell.",
    15: "**Native table — Table III of the report** (GSM8K accuracy), drawn by "
        "`grid()`. Rows match the report cell for cell.",
    7: "**Native table — Table II of the report** (configuration ladder), "
       "drawn by `grid()`. Rows match the report cell for cell, pool column "
       "included, because the pool sizes are what make the comparison fair.",
}


def read_deck():
    prs = Presentation(DECK)
    figs = json.load(open(MANI, encoding="utf-8"))
    fi = 0
    out = []
    for n, slide in enumerate(prs.slides, 1):
        rec = {"n": n, "title": "", "cite": "", "blocks": [], "figure": None,
               "extra": [], "table": None}
        boxes = []
        for sh in slide.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                rec["figure"] = figs[fi]
                fi += 1
            if sh.shape_type == MSO_SHAPE_TYPE.TABLE:
                rec["table"] = [[c.text for c in row.cells] for row in sh.table.rows]
            if (not sh.has_text_frame or not sh.text_frame.text.strip()
                    or sh.shape_type != MSO_SHAPE_TYPE.TEXT_BOX):
                continue
            boxes.append((sh.top / EMU, sh.left / EMU, sh))
        if n == 1:                             # the title slide has no chrome() and
            rec["title"] = "Title slide"       # no bullet structure to parse
            for top, left, sh in sorted(boxes):
                rec["extra"] += [x.text.strip() for x in sh.text_frame.paragraphs
                                 if x.text.strip()]
            rec["script"] = slide.notes_slide.notes_text_frame.text.strip()
            out.append(rec)
            continue
        for top, left, sh in sorted(boxes):
            txt = sh.text_frame.text.strip()
            if top < 0.5 and left < 1.0 and not rec["title"]:
                rec["title"] = txt
                continue
            if top > 6.8 and left < 1.0:
                rec["cite"] = txt
                continue
            if top > 6.8:
                continue                       # the slide number
            cur = None
            for p in sh.text_frame.paragraphs:
                if not p.text.strip():
                    continue
                is_head = False
                try:
                    is_head = bool(p.font.bold) and str(p.font.color.rgb) == ACCENT
                except Exception:
                    pass
                if is_head and p.level == 0:
                    cur = (p.text.strip(), [])
                    rec["blocks"].append(cur)
                elif cur is not None and p.level == 1:
                    cur[1].append(p.text.strip().lstrip("– ").strip())
                else:
                    rec["extra"].append(p.text.strip())
        rec["script"] = slide.notes_slide.notes_text_frame.text.strip()
        out.append(rec)
    return out


HEAD = """# Final Presentation Guide — build directions + full speaker script

**Course:** CSCI 6806 Graduate Capstone · **Group 3** · **Slot:** 10 minutes
**Deck:** {nslides} slides, `results/presentation/final-presentation.pptx` — no backup slides,
every figure **and every table** in the report appears and is explained on its own slide.

**Timing (word-count checked at 140 wpm, printed in each slide's notes):**

| Path | Time | What it is |
|---|---|---|
| **Full script** | **{full}** | every slide narrated in full |
| **Short path** (recommended) | **{short}** | slides {shortlist} delivered from the one-line `SHORT PATH` version in the notes |

The short path fits the 10-minute slot with ~45 s spare for handoffs and questions.
Nothing is deleted for it — the full text stays in the notes if a question sends you back.

> ⚠️ The requirement says *"the use of generative AI tools for content creation is prohibited."*
> Treat this as a **structured outline of your own results** and rewrite the sentences in your own
> voice before presenting. The figures and every number in them come from your committed runs.

---

## 1. Work split

Three presenters, ~3 minutes each. The boundaries fall on the natural section joints.

| Presenter | Slides | Full | Short path | Why this split |
|---|---|---|---|---|
| **Guoliang Liu** | {g_sl} (title, problem, data path, kernel, what we built) | {g_full} | {g_short} | You built the measurement harness, so the framing and the system description are yours |
| **Wenhui Kang** | {w_sl} (platform/Table I, design/Table II, reproduction, headline, cross-engine, tail) | {w_full} | {w_short} | You wrote the project and the report — you take the evaluation design and the main result |
| **Junpeng Huang** | {j_sl} (kernel sign flip, compute density, mechanism, accuracy/Table III, conclusion) | {j_full} | {j_short} | You built the LTR ranker and the baselines — you close the loop back to the scheduler |

**Screen share:** one person only (recommend Wenhui — you own the repo and can pull up code if asked).

**Handoff lines** (say them; they make the split obvious to the grader):
- Guoliang → Wenhui: *"Wenhui will now show the platform and how we made the comparison fair."*
- Wenhui → Junpeng: *"Junpeng will take the two conditions on when sparsity actually pays."*

---

## 2. The one analogy, used all the way through

Define it **once** on slide 2, then reuse the same words on 5 and 14. Never introduce a second metaphor.

| Real system | Analogy |
|---|---|
| A request | **a student** taking an exam |
| The KV cache | that student's **cheat sheet** |
| One KV block | **one page** of the cheat sheet |
| Output tokens | the **answers** the student writes |
| Re-reading the cache each token | before every answer, the student **re-reads the whole cheat sheet** |
| GPU memory | the **desk** — it only holds so many cheat sheets |
| Preemption | the desk is full → a student **leaves** and must come back later |
| Swap preemption | their cheat sheet is **kept in a locker** |
| Recompute preemption | their cheat sheet is **thrown away** — they must rewrite it |
| **C1** — 4-bit quantization | **smaller handwriting** → 2.83× more pages fit on the same desk |
| **C2** — attention sparsity | **skip re-reading the pages** that don't matter for this question |
| LTR scheduling | decides **which student sits down next** |
| Our layer | decides **how much desk space each student needs** |

---

## 3. Global build directions

**Colour theme** (rubric item 4). One accent, two neutrals, used consistently:
- Accent **teal `#0F766E`** — our contribution (C1, the cache layer, key numbers)
- Text **`#1A1A2E`**, muted **`#5A6672`**, panel fill **`#E8EEEE`**, background white
- In the charts the per-configuration colours are fixed and identical across every figure:
  B0 red `#9A3324` · B1 gold `#B8860B` · +C2 purple `#6A5ACD` · +C1 teal `#0F766E` · +C1+C2 green `#1B7A4B`
- Never use the accent for baseline/FCFS elements.

**Type:** titles 25 pt, main bullets 14.5–17 pt, sub-bullets 11.5–13.5 pt, citation line 10.5 pt.
Chart type is 15–16 pt at the placed size — **insert the PNGs at 100 %, do not stretch them**.

**Figures** (rubric item 5 — "original, not screenshots"). All nine are in
`results/figures/slides/`, regenerated at 300 dpi from your committed run summaries, and each maps
to a figure in the report so the deck and the paper can never disagree:

{figtable}

Regenerate any time with `.venv/bin/python latex_source/scripts/slide_figures.py`,
then rebuild the deck with `.venv/bin/python latex_source/scripts/build_deck.py`.
**Do not screenshot the report PDF** — that is exactly what the rubric penalizes.

**Animations** (required): one *Appear* or *Fade* per main bullet, in the click order listed per
slide below. No spins, no bounces, no sounds. On figure slides animate the figure in **after** the
bullets, so the audience hears the claim before seeing the chart.

**In-text citations** (required): every slide that makes a factual claim carries a citation line at
the bottom. Numbering matches the report's reference list:
`[1]` PagedAttention/vLLM · `[2]` Fu et al., LTR · `[3]` Kumar et al. · `[4]` TurboQuant ·
`[5]` KVmix · `[6]` BLASST · `[7]` HeadInfer · `[8]` MELL · `[9]` DistServe.

---

# 4. Slide by slide

Each section below is read straight out of the built deck: the bullets are the text on the slide,
and the script is the text in that slide's notes.

---
"""

TAIL = """
---

## 5. Rubric checklist — tick before you submit

- [ ] **Title slide** carries the project title, all three names, the university and the date
- [ ] **Equal contribution** — 5 / 6 / 5 slides, roughly 3 minutes each, named in the notes
- [ ] **Consistent colour theme** — one teal accent, fixed per-configuration chart colours
- [ ] **Original figures** — nine charts regenerated from run summaries, two native PowerPoint
      schematics, three native tables reproducing Tables I–III; no screenshots
- [ ] **Figure legibility** — 300 dpi, 15–16 pt chart type, inserted at 100 %
- [ ] **In-text citations** on every slide that makes a factual claim
- [ ] **Animations** — one per main bullet, figure last
- [ ] **Finishes inside 10 minutes** — short path is {short}; rehearse with a timer
- [ ] Exported to **PDF** as well as `.pptx` (the requirement asks for one PDF deck per group)

---

## 6. Likely questions, short answers

**"Isn't the four-bit configuration just getting more memory?"**
No — that is what the equal-byte design on slide 7 (Table II) controls for. Both pools occupy approximately the
same device memory; four-bit fits more *tokens* into those same bytes. And the `ctrl` row gives bf16
the same token capacity, so we can separate "more capacity helps" from "four-bit helps".

**"Why is end-to-end latency worse with compression alone?"**
Because compression trades decode work for memory. It cuts the queueing component of latency but
adds unpacking cost to every decode step. Slide 9 shows both halves; that is why we need C2.

**"Why does the same sparsity method help in one place and hurt in another?"**
Grouped-query attention. Four query heads share one execution tile, so a block can only be skipped
if all four agree. Our gate identifies ~42 % of blocks as negligible but the grouped kernel realizes
only ~18.5 %, which is below the break-even for the gate's own cost. Slide 12.

**"Is 100 GSM8K problems enough to claim no quality loss?"**
It is enough to claim *no measured degradation*, which is exactly what slide 15 says. It is not enough
to claim the system is more accurate — sparsity alone scores highest at 0.88, which we attribute to
noise rather than to reasoning gains — and we list the sample size as a limitation.

**"Why cap the KV block pool?"**
An uncapped RTX 3090 becomes compute-limited before the cache fills, so there would be no memory
pressure to relieve and the experiment could not test the hypothesis. The cap is stated in the
report as a design decision, not hidden. Slide 7.

**"Did you modify vLLM?"**
Not the installed package. C2 is installed by a `sitecustomize` module on `PYTHONPATH` that rewrites
the compiled decode kernel in memory inside vLLM's spawned EngineCore workers. It is reversible, and
a reviewer can diff it.

**"Would this help on an A100 or H100?"**
Compression would help whenever the server is memory-bound, which is common at high concurrency and
long context. Sparsity depends on the kernel: on a per-head kernel we measured a clear gain; on the
unified grouped kernel it costs. That is the conditional conclusion on the final slide, and replicating it
on a datacenter GPU is our stated next step.

---

## 7. Rehearsal plan

1. **Read-through, no slides.** Each person reads their own notes aloud once with a timer. If your
   section runs long, switch the marked slides to their `SHORT PATH` line rather than speeding up.
2. **With slides, separately.** Check every click lands on the right bullet and that the figure
   arrives last.
3. **Full run with handoffs.** Say the handoff lines out loud. Time the whole thing — target {short},
   hard ceiling 10:00.
4. **Q&A drill.** One person reads a question from §6 at random; whoever owns that slide answers.
5. **Final check.** Open the exported PDF on the machine you will present from and confirm the fonts
   and figures survived the export.
"""


def render():
    slides = read_deck()
    full = short = 0
    by = {}
    for s in slides:
        m = re.match(r"\[(\d+):(\d+) — (\w+)\]", s["script"])
        sec = int(m.group(1)) * 60 + int(m.group(2))
        x = re.search(r"SHORT PATH \((\d+):(\d+)\)", s["script"])
        sp = int(x.group(1)) * 60 + int(x.group(2)) if x else sec
        s["sec"], s["sp"], s["who"] = sec, sp, m.group(3)
        full += sec
        short += sp
        by.setdefault(m.group(3), [0, 0, []])
        by[m.group(3)][0] += sec
        by[m.group(3)][1] += sp
        by[m.group(3)][2].append(s["n"])

    def ts(v):
        return "%d:%02d" % (v // 60, v % 60)

    rows = ["| Slide figure | Report figure | Used on | Placed size |",
            "|---|---|---|---|"]
    for s in slides:
        if s["figure"]:
            f = s["figure"]
            rows.append("| `%s` | %s | Slide %d | %.2f × %.2f in |" % (
                f["figure"], REPORT_FIG.get(f["figure"], "—"), s["n"], f["w"], f["h"]))
    shortlist = [str(x["n"]) for x in slides if x["sp"] < x["sec"]]
    shortlist = ", ".join(shortlist[:-1]) + " and " + shortlist[-1]
    sl = {k: "%d–%d" % (v[2][0], v[2][-1]) for k, v in by.items()}
    body = [HEAD.format(
        full=ts(full), short=ts(short), figtable="\n".join(rows),
        nslides=len(slides), shortlist=shortlist,
        g_sl=sl["Guoliang"], w_sl=sl["Wenhui"], j_sl=sl["Junpeng"],
        g_full=ts(by["Guoliang"][0]), g_short=ts(by["Guoliang"][1]),
        w_full=ts(by["Wenhui"][0]), w_short=ts(by["Wenhui"][1]),
        j_full=ts(by["Junpeng"][0]), j_short=ts(by["Junpeng"][1]))]

    for s in slides:
        t = ["## Slide %d — %s" % (s["n"], s["title"]),
             "",
             "*%s* · %s%s" % (s["who"], ts(s["sec"]),
                              "  (short path %s)" % ts(s["sp"]) if s["sp"] < s["sec"] else ""),
             ""]
        if s["blocks"]:
            t += ["**Paste into PowerPoint:**", "```"]
            for head, subs in s["blocks"]:
                t.append(head)
                t += ["  - " + x for x in subs]
            t += ["```", ""]
        if s["extra"] and s["n"] == 1:
            t += ["**Paste into PowerPoint:**", "```"] + s["extra"] + ["```", ""]
        elif s["extra"]:
            t += ["**Other text on the slide:** " + " · ".join(s["extra"]), ""]
        if s["n"] in NATIVE:
            t += ["**Graphic:** " + NATIVE[s["n"]], ""]
        if s["table"]:
            hdr, *body_rows = s["table"]
            t += ["| " + " | ".join(hdr) + " |",
                  "|" + "---|" * len(hdr)]
            t += ["| " + " | ".join(r) + " |" for r in body_rows]
            t += [""]
        if s["figure"]:
            f = s["figure"]
            t += ["**Figure:** `results/figures/slides/%s` — %s in the report. "
                  "Placed at %.2f × %.2f in, x = %.2f, y = %.2f. Insert at 100 %%." % (
                      f["figure"], REPORT_FIG.get(f["figure"], "—"),
                      f["w"], f["h"], f["x"], f["y"]), ""]
        if len(s["blocks"]) > 1:
            t += ["**Click order:** " + " → ".join(
                "%d. %s" % (i + 1, h.split("  ")[0][:38])
                for i, (h, _) in enumerate(s["blocks"]))
                + (" → %d. figure" % (len(s["blocks"]) + 1) if s["figure"] else ""), ""]
        if s["cite"]:
            t += ["**Citation line (bottom of slide):** " + s["cite"], ""]
        t += ["**Say:**", ""]
        for para in s["script"].split("\n\n")[1:]:
            para = " ".join(para.split())
            if para.startswith("SHORT PATH"):
                t += ["> **" + para.split("—")[0].strip() + "** — "
                      + "—".join(para.split("—")[1:]).strip(), ""]
            else:
                t += ["> " + para, ""]
        t += ["---", ""]
        body.append("\n".join(t))

    body.append(TAIL.format(short=ts(short)))
    open(OUT, "w", encoding="utf-8").write("\n".join(body))
    return OUT, full, short


if __name__ == "__main__":
    path, f, s = render()
    print("wrote %s  (full %d:%02d, short path %d:%02d)" % (path, f // 60, f % 60, s // 60, s % 60))
